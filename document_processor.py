# document_processor.py

# Standard library imports
import os
import sys
import io
import time
import base64
import json
import hashlib
import logging
import asyncio
from datetime import datetime
from pathlib import Path
from typing import (
    List,
    Dict,
    Tuple,
    Optional,
    Generator,
    Callable,
    Type,
    Union,
    Any
)
from dataclasses import dataclass

# Third-party imports
# PDF Processing
import fitz  # PyMuPDF
from pptx import Presentation

# Image Processing
from PIL import Image, ImageEnhance
import imghdr
import numpy as np
import pandas as pd

# Machine Learning & AI
import torch
from transformers import CLIPProcessor, CLIPModel
from openai import OpenAI
from anthropic import Anthropic
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Vector Database
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.models import (
    Distance,
    VectorParams,
    PointStruct,
    OptimizersConfigDiff
)

# Web and API
import requests
from requests.auth import HTTPBasicAuth
from bs4 import BeautifulSoup
from tqdm import tqdm
import openpyxl

# Optional: For type hints in modern Python
from typing_extensions import TypedDict, Protocol

# deployment code
from clearml import Task

@dataclass
class ImageMetadata:
    source_path: str
    page_number: int
    image_number: int
    dimensions: Tuple[int, int]
    format: str
    dpi: Optional[Tuple[int, int]]
    file_size: int
    hash: str
    # New semantic fields
    image_caption: str
    text_in_image: str
    image_category: List[str]


class FileProcessor:
    def __init__(self, config):
        """Initialize FileProcessor with image processing configurations."""
        self.config = config
        self.image_extensions = self.config.SUPPORTED_IMAGE_FORMATS
        self.max_image_dimension = self.config.MAX_IMAGE_SIZE
        self.min_image_dimension = self.config.MIN_IMAGE_SIZE
        self.target_dpi = self.config.TARGET_IMAGE_DPI
        self.quality = self.config.IMAGE_QUALITY
        
        # Set up logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)


    def standardize_image(self, image: Image.Image) -> Image.Image:
        """Standardize image dimensions and quality."""
        try:
            # Get original dimensions
            width, height = image.size
            
            # Calculate aspect ratio
            aspect_ratio = width / height

            # Determine new dimensions while maintaining minimum resolution for text
            min_text_resolution = self.target_dpi  # DPI for text readability
            
            # Determine new dimensions
            if width > self.max_image_dimension or height > self.max_image_dimension:
                if width > height:
                    new_width = self.max_image_dimension
                    new_height = int(new_width / aspect_ratio)
                else:
                    new_height = self.max_image_dimension
                    new_width = int(new_height * aspect_ratio)
                
                # Resize image
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS, reducing_gap=3.0)  # Helps preserve detail

                enhancer = ImageEnhance.Sharpness(image)
                image = enhancer.enhance(1.2)

            return image
        except Exception as e:
            self.logger.error(f"Error standardizing image: {str(e)}")
            raise

    def get_image_metadata(self, image_data: bytes, source_path: str, page_number: int,
                           image_number: int) -> ImageMetadata:
        """Generate metadata for an image."""
        try:
            # Get basic image info
            image = Image.open(io.BytesIO(image_data))
            dimensions = image.size

            return ImageMetadata(
                source_path=source_path,
                page_number=page_number,
                image_number=image_number,
                dimensions=dimensions,
                format=imghdr.what(None, h=image_data) or 'png',
                dpi=image.info.get('dpi'),
                file_size=len(image_data),
                hash=hashlib.md5(image_data).hexdigest(),
                # Add default values for semantic fields
                image_caption='',
                text_in_image='',
                image_category=[]
            )

        except Exception as e:
            self.logger.error(f"Error generating image metadata: {str(e)}")
            # Return metadata with default values
            return ImageMetadata(
                source_path=source_path,
                page_number=page_number,
                image_number=image_number,
                dimensions=(0, 0),
                format='unknown',
                dpi=None,
                file_size=len(image_data),
                hash=hashlib.md5(image_data).hexdigest(),
                image_caption='',
                text_in_image='',
                image_category=[]
            )

    def process_pdf(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a single PDF file using PyMuPDF."""
        try:
            doc = fitz.open(file_path)

            for page_number in range(len(doc)):
                page = doc[page_number]

                # Get regular text FIRST
                text = page.get_text("text", flags=fitz.TEXT_PRESERVE_LIGATURES |
                                                   fitz.TEXT_PRESERVE_WHITESPACE)

                # If no text found, try blocks mode
                if not text.strip():
                    blocks = page.get_text("blocks")
                    if blocks:
                        text = "\n".join([block[4] for block in blocks if isinstance(block, tuple)])

                # First try to extract tables
                tables = self._extract_tables(page)

                if text.strip() or tables:
                    lines = text.strip().split('\n')
                    header = lines[0] if lines else f"Header for page {page_number + 1}"
                    content = '\n'.join(lines[1:]) if len(lines) > 1 else text

                    # Combine regular text and table content
                    final_content = self._combine_text_and_tables(content, tables)

                    images = self._extract_page_images(doc, page)
                    yield page_number, final_content.strip(), header, images
                else:
                    print(f"No content extracted from page {page_number + 1}")

                print(f"\n[DEBUG] Table extraction summary for page {page_number + 1}:")
                print(f"Number of tables found: {len(tables)}")
                for i, table in enumerate(tables):
                    print(f"\nTable {i + 1}:")
                    print(f"Rows: {len(table['data'])}")
                    print(f"Columns: {len(table['data'][0]) if table['data'] else 0}")

            doc.close()

        except Exception as e:
            self.logger.error(f"Error processing PDF {file_path}: {str(e)}")
            yield from []

    def _extract_tables(self, page: fitz.Page) -> List[Dict]:
        """Extract tables from the page using different methods."""
        tables = []
        try:
            # Use tabular text extraction
            tabular = page.get_text("tables")
            if tabular:
                for table in tabular:
                    processed_table = self._process_raw_table(table)
                    if processed_table:
                        tables.append(processed_table)

            return tables

        except Exception as e:
            self.logger.error(f"Error extracting tables: {str(e)}")
            return []

    def _process_table_block(self, block: Dict) -> Optional[Dict]:
        """Process a table block from structured text extraction."""
        try:
            rows = []
            current_row = []
            current_y = None

            for line in block.get("lines", []):
                y = line["bbox"][1]  # y-coordinate of the line

                # New row detection
                if current_y is None or abs(y - current_y) > 5:  # threshold for new row
                    if current_row:
                        rows.append(current_row)
                    current_row = []
                    current_y = y

                # Process spans in the line
                text = " ".join(span["text"] for span in line.get("spans", []))
                current_row.append(text)

            # Add last row
            if current_row:
                rows.append(current_row)

            return {
                "type": "table",
                "data": rows,
                "bbox": block["bbox"]
            }

        except Exception as e:
            self.logger.error(f"Error processing table block: {str(e)}")
            return None

    def _process_raw_table(self, table: List[List]) -> Optional[Dict]:
        """Process a raw table from tabular text extraction."""
        try:
            # Minimum requirements for a valid table
            MIN_ROWS = 2
            MIN_COLS = 2

            # Initial dimension check
            if not table or len(table) < MIN_ROWS:
                return None

            # Clean and validate rows
            processed_rows = []
            for row in table:
                cleaned_row = [str(cell).strip() for cell in row if cell is not None]
                if len(cleaned_row) >= MIN_COLS and any(cleaned_row):  # Ensure row has content
                    processed_rows.append(cleaned_row)

            # Final validation
            if (len(processed_rows) >= MIN_ROWS and
                all(len(row) >= MIN_COLS for row in processed_rows)):
                return {
                    "type": "table",
                    "data": processed_rows
                }
            return None

        except Exception as e:
            self.logger.error(f"Error processing raw table: {str(e)}")
            return None

    def _combine_text_and_tables(self, text: str, tables: List[Dict]) -> str:
        """Combine regular text and table content."""
        result = []

        if text:
            result.append(text)

        for table in tables:
            result.append("\n=== Table Start ===\n")

            if table["data"]:
                # Calculate column widths
                col_widths = [0] * len(table["data"][0])
                for row in table["data"]:
                    for i, cell in enumerate(row):
                        col_widths[i] = max(col_widths[i], len(str(cell)))

                # Format table with proper alignment
                for row in table["data"]:
                    formatted_row = " | ".join(
                        str(cell).ljust(col_widths[i])
                        for i, cell in enumerate(row)
                    )
                    result.append(formatted_row)

            result.append("=== Table End ===\n")

        return "\n".join(result)

    def _extract_page_images(self, doc: fitz.Document, page: fitz.Page) -> List[Dict]:
        """Extract images from PDF page using direct pixel access."""
        images = []
        try:
            image_list = page.get_images()
            print(f"\n[DEBUG] Page {page.number + 1}: Found {len(image_list)} images")
            print("[DEBUG] Raw image_list structure:")
            print(f"{image_list}")

            for img_index, img_info in enumerate(image_list):
                try:
                    print(f"\n{'=' * 50}")
                    print(f"[DEBUG] Processing image {img_index + 1}/{len(image_list)}")

                    if len(img_info) < 1:
                        continue

                    xref = img_info[0]

                    # Print raw properties for debugging
                    print("\n[DEBUG] Raw Image Properties (img_info):")
                    properties = ['xref', 'smask', 'width', 'height', 'bpc', 'colorspace',
                                  'alt. colorspace', 'filter', 'interpolate']
                    for i, prop in enumerate(properties):
                        if i < len(img_info):
                            print(f"  {prop}: {img_info[i]}")

                    # Try alternative extraction method
                    try:
                        # Get the image mask if it exists
                        mask_xref = img_info[1] if len(img_info) > 1 else None

                        # Extract image using pixmap
                        pix = fitz.Pixmap(doc, xref)

                        # Handle color space conversion if needed
                        if pix.colorspace.n >= 4:  # CMYK or other
                            pix = fitz.Pixmap(fitz.csRGB, pix)

                        # If there's a mask, apply it
                        if mask_xref and mask_xref != 0:
                            try:
                                mask_pix = fitz.Pixmap(doc, mask_xref)
                                pix = fitz.Pixmap(pix, mask_pix)
                            except Exception as e:
                                print(f"[DEBUG] Note: Error applying mask: {str(e)}")

                        # Convert to PNG bytes
                        img_bytes = pix.tobytes("png")

                        print("\n[DEBUG] Extracted Image Properties:")
                        print(f"  Width: {pix.width}")
                        print(f"  Height: {pix.height}")
                        print(f"  Color Space: {pix.colorspace.n}")
                        print(f"  Alpha Channel: {pix.alpha}")
                        print(f"  Size: {len(img_bytes)} bytes")

                        image_dict = {
                            'image_data': img_bytes,
                            'format': 'png',
                            'width': pix.width,
                            'height': pix.height,
                            'colorspace': pix.colorspace.n,
                            'has_alpha': pix.alpha,
                            'xref': xref,
                            'related_text': f"Image {img_index + 1} from page {page.number + 1}"
                        }

                        # Clean up
                        pix = None

                        images.append(image_dict)
                        print(f"\n[DEBUG] ✓ Successfully processed image {img_index + 1}")
                        print(f"Final image size: {len(image_dict['image_data'])} bytes")

                    except Exception as e:
                        print(f"[DEBUG] ❌ Error in alternative extraction: {str(e)}")
                        # Fall back to original extract_image method
                        base_image = doc.extract_image(xref)
                        if base_image and 'image' in base_image:
                            print("[DEBUG] Using fallback extraction method")
                            image_dict = {
                                'image_data': base_image['image'],
                                'format': base_image.get('ext', 'png'),
                                'width': base_image.get('width'),
                                'height': base_image.get('height'),
                                'xref': xref,
                                'related_text': f"Image {img_index + 1} from page {page.number + 1}"
                            }
                            images.append(image_dict)

                    print(f"{'=' * 50}\n")

                except Exception as e:
                    print(f"[DEBUG] ❌ Error processing image {img_index + 1}: {str(e)}")
                    self.logger.error(f"Error extracting image {img_index}: {str(e)}")
                    continue

            print(f"\n[DEBUG] Page Processing Summary:")
            print(f"  - Total images found: {len(image_list)}")
            print(f"  - Successfully processed: {len(images)}")
            print(f"  - Failed: {len(image_list) - len(images)}")

            return images

        except Exception as e:
            print(f"[DEBUG] ❌ Fatal error processing page: {str(e)}")
            self.logger.error(f"Error processing page images: {str(e)}")
            return images

    def _enhance_image_quality(self, image: Image.Image) -> Image.Image:
        """Apply quality enhancements to improve clarity."""
        try:
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.5)

            # Enhance contrast slightly
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)

            # Enhance color
            enhancer = ImageEnhance.Color(image)
            image = enhancer.enhance(1.1)

            return image

        except Exception as e:
            self.logger.error(f"Error enhancing image: {str(e)}")
            return image

    def process_pptx(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a PPTX file with enhanced image handling."""
        try:
            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)
            processed_hashes = set()  # Track duplicate images

            self.logger.info(f"Processing PPTX: {file_path} with {total_slides} slides")

            for slide_number, slide in enumerate(presentation.slides, 1):
                self.logger.info(f"Processing slide {slide_number}/{total_slides}")
                slide_content = self._extract_slide_content(slide)
                images = self._extract_slide_images(slide)

                # Filter out duplicate images
                unique_images = []
                for img in images:
                    if img['hash'] not in processed_hashes:
                        processed_hashes.add(img['hash'])
                        unique_images.append(img)

                # Log content and images found
                self.logger.info(f"Slide {slide_number}: Found {len(unique_images)} unique images")
                self.logger.info(f"Slide {slide_number}: Text length: {len(slide_content['text'])}")

                if slide_content['text'] or slide_content['title'] or unique_images:
                    yield slide_number, slide_content['text'], slide_content['title'], unique_images

        except Exception as e:
            self.logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            yield from []


    def _extract_slide_content(self, slide) -> Dict[str, str]:
        """Helper method to extract content from a slide."""
        content = {
            'title': '',
            'text': []
        }
        
        if slide.shapes.title:
            content['title'] = slide.shapes.title.text.strip()
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if text != content['title']:
                    content['text'].append(text)
                    
        return {
            'title': content['title'],
            'text': '\n'.join(content['text'])
        }

    def _extract_slide_images(self, slide) -> List[Dict]:
        """Extract original images from PowerPoint slide with enhanced validation and metadata."""
        images = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        if not image_bytes:
                            continue

                        # Validate image data
                        try:
                            image = Image.open(io.BytesIO(image_bytes))
                            format_type = image.format.lower()
                            dimensions = image.size

                            # Basic image validation
                            if dimensions[0] < 10 or dimensions[1] < 10:
                                self.logger.warning(f"Image too small on slide {slide.slide_id}: {dimensions}")
                                continue

                            # Get surrounding text and title
                            slide_title = slide.shapes.title.text.strip() if slide.shapes.title else ""

                            # Calculate image hash for duplicate detection
                            image_hash = hashlib.md5(image_bytes).hexdigest()

                            images.append({
                                'image_data': image_bytes,
                                'format': format_type,
                                'dimensions': dimensions,
                                'hash': image_hash,
                                'slide_number': slide.slide_id,
                                'related_text': f"{slide_title}\nImage from slide {slide.slide_id}",
                                'alt_text': shape.alt_text if hasattr(shape, 'alt_text') else ""
                            })

                        except (IOError, OSError) as img_err:
                            self.logger.error(f"Invalid image data on slide {slide.slide_id}: {str(img_err)}")
                            continue

                    except AttributeError as attr_err:
                        self.logger.error(
                            f"Error accessing image attributes on slide {slide.slide_id}: {str(attr_err)}")
                        continue

        except Exception as e:
            self.logger.error(f"Error processing images in slide {slide.slide_id}: {str(e)}")

        return images

    def process_xlsx(self, file_path: str) -> Generator[Tuple[int, str, str, list], None, None]:
        """
        Process xlsx files and yield data in format: (sheet_number, content, header, images)
        content: contains vectorizable text (filepath, filename, sheet names, headers)
        """
        try:
            print(f"Starting to process Excel file: {file_path}")

            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            file_name = Path(file_path).name
            file_path = str(Path(file_path))

            # Process each sheet
            for sheet_number, sheet_name in enumerate(excel_file.sheet_names):
                print(f"Processing sheet {sheet_number + 1}: {sheet_name}")

                df = pd.read_excel(excel_file, sheet_name=sheet_name)

                # Get column names (headers)
                headers = list(df.columns)

                # Construct content for vectorization
                content = f"File: {file_name}\n"
                content += f"Path: {file_path}\n"
                content += f"Sheet: {sheet_name}\n"
                content += f"Headers: {', '.join(headers)}"

                print(f"Successfully processed sheet {sheet_name} with {len(headers)} columns")

                # Yield in format compatible with existing processor
                # (sheet_number, content, header, images)
                yield (sheet_number, content, sheet_name, [])

            print(f"Successfully completed processing Excel file: {file_name}")
            print(f"Total sheets processed: {len(excel_file.sheet_names)}")

        except Exception as e:
            print(f"Error processing Excel file {file_path}: {str(e)}")
            yield (0, "", "", [])


    def process_document(self, file_path: str) -> Generator[tuple, None, None]:
        """Process any supported document type."""
        file_extension = os.path.splitext(file_path)[1].lower()

        processors = {
            '.pdf': self.process_pdf,
            '.pptx': self.process_pptx,
            '.ppt': self.process_pptx,
            '.xlsx': self.process_xlsx
        }
        
        processor = processors.get(file_extension)
        if processor:
            yield from processor(file_path)
        else:
            self.logger.warning(f"Unsupported file type: {file_extension}")
            yield from []

    def save_image(self, image_data: bytes, source_file: str, page_number: int,
                  image_number: int, storage_path: str, original_format: str = 'png') -> str:
        """Save original image bytes without modification."""
        try:
            print("\n\nSource File : ", source_file)
            # Create directory if it doesn't exist
            storage_path = self.config.IMAGE_STORAGE_PATH
            os.makedirs(storage_path, exist_ok=True)

            # Generate filename based on source and location
            base_name = os.path.splitext(os.path.basename(source_file))[0]
            print("\nbase name : ", base_name)
            filename = f"{base_name}_page{page_number}_img{image_number}.{original_format}"
            file_path = os.path.join(storage_path, filename)

            # Write original bytes directly to file
            with open(file_path, 'wb') as f:
                f.write(image_data)

            return str(file_path)

        except Exception as e:
            self.logger.error(f"Error saving image: {str(e)}")
            raise


class ConfluenceProcessor:
    def __init__(self, user_name: str, api_token: str, base_url: str):
        self.api_token = api_token
        self.user_name = user_name
        self.base_url = base_url.rstrip('/')

    def process_confluence(self, page_id: str) -> Generator[tuple, None, None]:
        """Fetch and process content from a Confluence page."""
        confluence_url = f'{self.base_url}/rest/api/content/{page_id}?expand=body.storage'
        try:
            response = requests.get(confluence_url, auth=HTTPBasicAuth(self.user_name, self.api_token))
            if response.status_code == 200:
                data = response.json()
                page_content = data['body']['storage']['value']
                soup = BeautifulSoup(page_content, 'html.parser')
                text = soup.get_text()

                images = []
                for img in soup.find_all('ac:image'):
                    attachment = img.find('ri:attachment')
                    if attachment:
                        filename = attachment['ri:filename']
                        img_url = f'{self.base_url}/download/attachments/{page_id}/{filename}'
                        images.append(img_url)

                yield 0, text.strip(), data['title'], images
            else:
                print(f"Failed to retrieve Confluence page: {response.status_code} - {response.text}")
        except Exception as e:
            print(f"Error processing Confluence page: {str(e)}")
            yield from []


class VectorStore:
    def __init__(self, config):
        # self.qdrant_client = QdrantClient(
        #     host="localhost",  # Local server host
        #     port=6333  # Default Qdrant server port
        # )

        # Update Qdrant initialization to use cloud configuration
        self.qdrant_client = QdrantClient(
            url=config.QDRANT_URL,
            api_key=config.QDRANT_API_KEY
        )

        self.collection_name = config.COLLECTION_NAME
        self.BATCH_SIZE = config.BATCH_SIZE
        self.image_collection_name = f"{config.COLLECTION_NAME}_images"
        self.logger = logging.getLogger(__name__)
        self._setup_collections()

    def recover_collection(self, collection_name: str):
        """Attempt to recover a potentially corrupted collection."""
        try:
            # Force reindexing
            self.qdrant_client.update_collection(
                collection_name=collection_name,
                optimizers_config=models.OptimizersConfigDiff(
                    indexing_threshold=1
                )
            )
            
            # Wait for recovery
            MAX_WAIT = 30  # seconds
            start_time = time.time()
            while time.time() - start_time < MAX_WAIT:
                collection_info = self.qdrant_client.get_collection(collection_name)
                if collection_info.status == "green":
                    return True
                time.sleep(2)
                
            return False
            
        except Exception as e:
            self.logger.error(f"Error recovering collection {collection_name}: {str(e)}")
            return False

    # Add this method to use recovery when needed
    def ensure_collection_health(self, collection_name: str) -> bool:
        """Ensure collection is healthy, attempt recovery if needed."""
        try:
            collection_info = self.qdrant_client.get_collection(collection_name)
            if collection_info.status != "green":
                self.logger.warning(f"Collection {collection_name} needs recovery")
                return self.recover_collection(collection_name)
            return True
        except Exception as e:
            self.logger.error(f"Error checking collection health: {str(e)}")
            return False

    def _setup_collections(self):
        """Setup Qdrant collections for both text and images if they don't exist."""
        try:
            collections = self.qdrant_client.get_collections()
            collection_names = [c.name for c in collections.collections]

            # Text collection setup
            if self.collection_name not in collection_names:
                self._create_text_collection()
            else:
                # Verify and recover existing collection if needed
                if not self.ensure_collection_health(self.collection_name):
                    raise Exception(f"Failed to recover text collection {self.collection_name}")

            # Image collection setup
            if self.image_collection_name not in collection_names:
                self._create_image_collection()
            
            else:
                # Verify and recover existing image collection if needed
                if not self.ensure_collection_health(self.image_collection_name):
                    raise Exception(f"Failed to recover image collection {self.image_collection_name}")
        
        except Exception as e:
            self.logger.error(f"Error setting up collections: {str(e)}")
            raise

    def _create_text_collection(self):
        """Create collection for text vectors."""
        try:
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=models.VectorParams(
                    size=1536,  # OpenAI embedding size
                    distance=models.Distance.COSINE
                ),
                optimizers_config=models.OptimizersConfigDiff(
                    indexing_threshold=1,  # Force immediate indexing
                    default_segment_number=1  # Use single segment for small collections
                ),
                # Add timeout and wait parameters
                timeout=60
            )

            # Verify collection creation and wait for it to be ready
            MAX_WAIT = 30  # seconds
            start_time = time.time()
            while time.time() - start_time < MAX_WAIT:
                collection_info = self.qdrant_client.get_collection(self.collection_name)
                if (collection_info and 
                    collection_info.status == "green" and 
                    collection_info.optimizer_status == "ok"):
                    return
                time.sleep(2)
            
        except Exception as e:
            self.logger.error(f"Error creating text collection: {str(e)}")
            raise

    def _create_image_collection(self):
        """Create collection for image vectors."""
        try:
            self.qdrant_client.create_collection(
                collection_name=self.image_collection_name,
                vectors_config=models.VectorParams(
                    size=512,  # CLIP embedding size
                    distance=models.Distance.COSINE
                ),
                optimizers_config=models.OptimizersConfigDiff(
                    indexing_threshold=1,  # Force immediate indexing
                    default_segment_number=1  # Use single segment for small collections
                ),
                timeout=60
            )

            # Verify collection creation and wait for it to be ready
            MAX_WAIT = 30  # seconds
            start_time = time.time()
            while time.time() - start_time < MAX_WAIT:
                collection_info = self.qdrant_client.get_collection(self.image_collection_name)
                if (collection_info and 
                    collection_info.status == "green" and 
                    collection_info.optimizer_status == "ok"):
                    return
                time.sleep(2)
                
            raise Exception("Collection creation verification timeout")
            
        except Exception as e:
            self.logger.error(f"Error creating image collection: {str(e)}")
            raise

    def store_text_vectors(self, vectors: List[List[float]], chunks: List[Dict]):
        """Store text vectors and metadata in Qdrant."""
        try:
            for i in range(0, len(vectors), self.BATCH_SIZE):
                batch_vectors = vectors[i:i + self.BATCH_SIZE]
                batch_chunks = chunks[i:i + self.BATCH_SIZE]
                points = []

                for chunk, vector in zip(batch_chunks, batch_vectors):

                    # Determine content type based on filename
                    filename = chunk['metadata']['filename']
                    content_type = 'excel' if filename.endswith(('.xlsx', '.xls', '.csv')) else 'text'

                    point = models.PointStruct(
                        id=self._generate_point_id(chunk['metadata']),
                        vector=vector,
                        payload={
                            'text': chunk['text'],
                            'filename': filename,
                            'page_number': chunk['metadata']['page_number'],
                            'chunk_number': chunk['metadata']['chunk_number'],
                            'page_header': chunk['metadata']['page_header'],
                            'timestamp': datetime.now().isoformat(),
                            'type': content_type  # Dynamic type based on file extension
                        }
                    )
                    points.append(point)

                self.qdrant_client.upsert(
                    collection_name=self.collection_name,
                    points=points
                )

                # Verify points were added
                collection_info = self.qdrant_client.get_collection(self.collection_name)
                if collection_info.vectors_count is None:
                    # Force refresh collection info
                    self.qdrant_client.update_collection(
                        collection_name=self.collection_name,
                        optimizers_config=models.OptimizersConfigDiff(
                            indexing_threshold=1  # Force immediate indexing
                        )
                    )
            
        except Exception as e:
            self.logger.error(f"Error storing text vectors: {str(e)}")
            raise

    def store_image_vectors(self, vectors: List[List[float]], image_metadata: List[ImageMetadata],
                            related_text: Optional[str] = None):
        """Store image vectors and metadata in Qdrant."""
        try:
            for i in range(0, len(vectors), self.BATCH_SIZE):
                batch_vectors = vectors[i:i + self.BATCH_SIZE]
                batch_metadata = image_metadata[i:i + self.BATCH_SIZE]
                points = []
                for metadata, vector in zip(batch_metadata, batch_vectors):
                    payload = {
                        'source_path': metadata.source_path,
                        'page_number': metadata.page_number,
                        'image_number': metadata.image_number,
                        'dimensions': list(metadata.dimensions),
                        'format': metadata.format,
                        'dpi': list(metadata.dpi) if metadata.dpi else None,
                        'file_size': metadata.file_size,
                        'hash': metadata.hash,
                        'image_caption': metadata.image_caption,
                        'text_in_image': metadata.text_in_image,
                        'image_category': metadata.image_category,
                        'related_text': related_text,
                        'timestamp': datetime.now().isoformat(),
                        'type': 'image'
                    }

                    # Log metadata for each image
                    self.logger.info(f"Storing image metadata:\n"
                                     f"Source: {payload['source_path']}\n"
                                     f"Page: {payload['page_number']}\n"
                                     f"Image: {payload['image_number']}\n"
                                     f"Caption: {payload['image_caption']}\n"
                                     f"Categories: {payload['image_category']}\n"
                                     f"Text: {payload['text_in_image']}")

                    points.append(models.PointStruct(
                        id=self._generate_image_point_id(metadata),
                        vector=vector,
                        payload=payload
                    ))

                self.qdrant_client.upsert(
                    collection_name=self.image_collection_name,
                    points=points
                )

                collection_info = self.qdrant_client.get_collection(self.image_collection_name)
                if collection_info.vectors_count is None:
                    self.qdrant_client.update_collection(
                        collection_name=self.image_collection_name,
                        optimizers_config=models.OptimizersConfigDiff(
                            indexing_threshold=1
                        )
                    )

                self.logger.info(f"Successfully stored batch of {len(points)} image vectors")

        except Exception as e:
            self.logger.error(f"Error storing image vectors: {str(e)}")
            raise

    def check_collection_status(self):
        """Check and wait for collection to be ready."""
        try:
            MAX_RETRIES = 5
            WAIT_TIME = 2  # seconds
            
            for attempt in range(MAX_RETRIES):
                collection_info = self.qdrant_client.get_collection(
                    collection_name=self.collection_name
                )
                
                if (collection_info.status == "green" and 
                    collection_info.optimizer_status == "ok"):

                    # Force indexing if vectors_count is None
                    if collection_info.vectors_count is None:
                        self.qdrant_client.update_collection(
                            collection_name=self.collection_name,
                            optimizers_config=models.OptimizersConfigDiff(
                                indexing_threshold=1,  # Force immediate indexing
                                default_segment_number=1  # Use single segment for small collections
                            )
                        )
                        # Wait for indexing to complete
                        time.sleep(WAIT_TIME)
                        continue
                    return True
                    
                if attempt < MAX_RETRIES - 1:
                    time.sleep(WAIT_TIME)
                    
            # Final check
            collection_info = self.qdrant_client.get_collection(
                collection_name=self.collection_name
            )
            return (collection_info.status == "green" and 
                collection_info.optimizer_status == "ok")
            
        except Exception as e:
            self.logger.error(f"Error checking collection status: {str(e)}")
            return False

    def _generate_point_id(self, metadata: Dict) -> int:
        """Generate a unique ID for text points."""
        id_string = f"{metadata['filename']}_{metadata['page_number']}_{metadata['chunk_number']}"
        return abs(hash(id_string)) % (2**63)

    def _generate_image_point_id(self, metadata: ImageMetadata) -> int:
        """Generate a unique ID for image points."""
        id_string = f"{metadata.source_path}_{metadata.page_number}_{metadata.image_number}_{metadata.hash}"
        return abs(hash(id_string)) % (2**63)

    def get_processed_files(self) -> set:
        """Get set of already processed files from both collections."""
        try:
            # Get text files
            text_response = self.qdrant_client.scroll(
                collection_name=self.collection_name,
                limit=10000,
                with_payload=['filename'],
                with_vectors=False
            )
            text_files = {point.payload['filename'] for point in text_response[0]}

            # Get image files
            image_response = self.qdrant_client.scroll(
                collection_name=self.image_collection_name,
                limit=10000,
                with_payload=['source_path'],
                with_vectors=False
            )
            image_files = {os.path.basename(point.payload['source_path']) 
                          for point in image_response[0]}

            return text_files.union(image_files)
        except Exception as e:
            self.logger.error(f"Error retrieving processed files: {str(e)}")
            return set()

    def search_similar(self, query_vector: List[float], limit: int = 10,
                       collection_name: Optional[str] = None) -> List[Dict]:
        """Search for similar vectors in specified collection."""
        try:
            collection = collection_name or self.collection_name
            response = self.qdrant_client.search(
                collection_name=collection,
                query_vector=query_vector,
                limit=limit
            )
            return [
                {
                    'payload': point.payload,
                    'score': point.score
                }
                for point in response
            ]
        except Exception as e:
            self.logger.error(f"Error searching vectors: {str(e)}")
            return []


class DocumentProcessor:
    def __init__(self, config, stats):

        # Initialize ClearML
        self.task = Task.init(
            project_name=config.CLEARML_PROJECT,
            task_name=config.CLEARML_TASK
        )

        # Initialize different counters for different metrics
        self.metrics_iterations = {
            'chunks': 0,
            'images': 0,
            'pages': 0
        }

        # Log configuration parameters
        self.task.connect(config)

        """Initialize document processor with all necessary components."""
        self.config = config
        self.stats = stats
        self.openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
        self.file_processor = FileProcessor(config)
        self.confluence_processor = ConfluenceProcessor(config.CONFLUENCE_USERNAME, config.CONFLUENCE_API_TOKEN, config.CONFLUENCE_BASE_URL)
        self.vector_store = VectorStore(config)
        self.logger = logging.getLogger(__name__)
        self.chunk_count = 0

        # Create image storage directory
        os.makedirs(config.IMAGE_STORAGE_PATH, exist_ok=True)

        # Initialize CLIP model for image processing
        self.device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
        self.clip_model = CLIPModel.from_pretrained(config.CLIP_MODEL).to(self.device)
        self.clip_processor = CLIPProcessor.from_pretrained(config.CLIP_PROCESSOR)
        self.semantic_analyzer = ImageSemanticAnalyzer(
            config.ANTHROPIC_API_KEY,
            config.SUPPORTED_IMAGE_FORMATS
        )

        # Text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ".", " ", ""]
        )

    def _get_image_embedding(self, image_data: bytes) -> List[float]:
        """Generate embedding for an image using CLIP."""
        try:
            # Convert bytes to PIL Image for CLIP processing only
            image = Image.open(io.BytesIO(image_data))

            # Convert to RGB only for embedding generation
            if image.mode != 'RGB':
                image = image.convert('RGB')

            inputs = self.clip_processor(images=image, return_tensors="pt")
            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**inputs)
            return image_features.squeeze().tolist()
        except Exception as e:
            self.logger.error(f"Error generating image embedding: {str(e)}")
            raise

    def create_chunks(self, text: str, metadata: Dict) -> List[Dict]:
        """Create chunks from text with metadata."""
        try:
            chunks = self.text_splitter.split_text(text)
            return [
                {
                    'text': chunk,
                    'metadata': {
                        'filename': metadata['filename'],
                        'page_number': metadata['page_number'],
                        'chunk_number': i + 1,
                        'page_header': metadata['page_header']
                    }
                }
                for i, chunk in enumerate(chunks)
            ]
        except Exception as e:
            self.logger.error(f"Error creating chunks: {str(e)}")
            return []

    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Get embeddings for texts using OpenAI."""
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-ada-002",
                input=texts
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            self.logger.error(f"Error generating text embeddings: {str(e)}")
            raise

    def process_document_page(self, file_path: str, page_number: int, text: str, header: str, images: List[Dict]) -> None:
        """Process a single page of a document."""
        chunks_processed = 0
        try:
            # Log page processing in clearML
            self.metrics_iterations['pages'] += 1
            self.task.get_logger().report_text(
                f"Processing page {page_number} of {file_path}"
            )

            # Process text content
            chunks = self.create_chunks(text, {
                'filename': os.path.basename(file_path),
                'page_number': page_number,
                'page_header': header
            })

            if chunks:
                text_embeddings = self.get_embeddings([c['text'] for c in chunks])
                self.vector_store.store_text_vectors(text_embeddings, chunks)
                chunks_processed = len(chunks)
                self.stats.total_chunks += chunks_processed

            # ClearML Log - Process chunks
            self.metrics_iterations['chunks'] += 1
            self.task.get_logger().report_scalar(
                "processing",
                "chunks_processed",
                iteration=self.metrics_iterations['chunks'],
                value=chunks_processed
            )

            # Process images
            for img_index, img_data in enumerate(images):
                try:
                    # Get semantic analysis
                    semantic_info = self.semantic_analyzer.analyze_image(img_data['image_data'])

                    # Save original image
                    stored_image_path = self.file_processor.save_image(
                        img_data['image_data'],
                        file_path,
                        page_number,
                        img_index,
                        self.config.IMAGE_STORAGE_PATH,
                        img_data['format']
                    )

                    metadata = ImageMetadata(
                        source_path=stored_image_path,
                        page_number=page_number,
                        image_number=img_index,
                        dimensions=(img_data.get('width', 0), img_data.get('height', 0)),
                        format=img_data['format'],
                        dpi=None,  # Original DPI if available
                        file_size=len(img_data['image_data']),
                        hash=hashlib.md5(img_data['image_data']).hexdigest(),
                        image_caption=semantic_info['image_caption'],
                        text_in_image=semantic_info['text_in_image'],
                        image_category=semantic_info['image_category']
                    )

                    embedding = self._get_image_embedding(img_data['image_data'])
                    self.vector_store.store_image_vectors(
                        [embedding],
                        [metadata],
                        img_data.get('related_text')
                    )

                    # ClearMl - Process images
                    if images:
                        self.metrics_iterations['images'] += 1
                        self.task.get_logger().report_scalar(
                            "processing",
                            "images_processed",
                            iteration=self.metrics_iterations['images'],
                            value=len(images)
                        )

                except Exception as e:
                    self.logger.error(f"Error processing image {img_index}: {str(e)}")
                    self.stats.errors.append(f"Error processing image {img_index}: {str(e)}")

            return chunks_processed

        except Exception as e:
            self.logger.error(f"Error processing page {page_number} of {file_path}: {str(e)}")
            self.stats.errors.append(f"Error processing page {page_number} of {file_path}: {str(e)}")
            return chunks_processed

    def final_processing(self) -> int:
        """Main processing function."""
        processed_files = self.vector_store.get_processed_files()
        files_processed = 0
        batch_processor = BatchProcessor(self.config.BATCH_SIZE)

        try:
            if not self.vector_store.check_collection_status():
                raise DocumentProcessingException("Collection not ready for processing")

            # Process PDF and PPT files
            for directory, extensions in [
                (self.config.PDF_DIRECTORY, ['.pdf']),
                (self.config.PPT_DIRECTORY, ['.ppt', '.pptx']),
                (self.config.EXCEL_DIRECTORY, ['.xlsx'])
            ]:
                # Fix #3: Check if directory exists and is valid
                if not os.path.exists(directory) or not os.path.isdir(directory):
                    self.logger.info(f"Skipping non-existent directory: {directory}")
                    continue

                # Get list of valid files in directory
                files = [f for f in os.listdir(directory)
                         if any(f.lower().endswith(ext) for ext in extensions)]

                # If directory is empty, skip to next directory
                if not files:
                    self.logger.info(f"No files found in directory: {directory}")
                    continue

                # Only update stats for non-empty directories
                self.stats.total_documents += len(files)

                for filename in tqdm(files, desc=f"Processing {os.path.basename(directory)}"):
                    file_key = f"{os.path.basename(directory)}_{filename}"
                    if file_key in processed_files:
                        self.logger.debug(f"Skipping already processed file: {file_key}")
                        continue

                    file_path = os.path.join(directory, filename)
                    batch_processor.add_item((file_path, filename))
                    files_processed += 1

                    if batch_processor.is_ready():
                        self._process_batch(batch_processor.current_batch)
                        batch_processor.current_batch = []

                # Fix #2: Process and clear any remaining items for this directory
                if batch_processor.current_batch:
                    self._process_batch(batch_processor.current_batch)
                    batch_processor.current_batch = []  # Clear before next directory

            # Process Confluence content
            self._process_confluence_content()

            # Fix #1: Don't override the chunk count
            return files_processed

        except Exception as e:
            self.logger.error(f"Error in final processing: {str(e)}")
            self.stats.errors.append(str(e))
            return files_processed

    def _process_batch(self, batch_items: List[Tuple[str, str]]) -> None:
        for file_path, filename in batch_items:
            try:
                for page_number, text, header, images in self.file_processor.process_document(file_path):
                    if text:
                        chunks_processed = self.process_document_page(file_path, page_number, text, header, images)
                        self.stats.total_pages += 1
                        self.stats.total_images += len(images)
                        self.chunk_count += chunks_processed
                        
            except Exception as e:
                self.logger.error(f"Error processing file {filename}: {str(e)}")
                self.stats.errors.append(f"Error processing file {filename}: {str(e)}")

    def _process_confluence_content(self):
        """Process Confluence content including images."""
        try:
            self.logger.info("Starting Confluence content processing...")
            for _, text, title, images in self.confluence_processor.process_confluence(
                self.config.CONFLUENCE_PAGE_ID
            ):
                if text:
                    self.logger.info("In Confluence content... processing text")
                    chunks = self.create_chunks(text, {
                        'filename': f"confluence_{self.config.CONFLUENCE_PAGE_ID}",
                        'page_number': 0,
                        'page_header': title
                    })
                    
                    if chunks:
                        text_embeddings = self.get_embeddings([c['text'] for c in chunks])
                        self.vector_store.store_text_vectors(text_embeddings, chunks)
                        self.stats.total_chunks += len(chunks)  # Update chunks count

                if images:
                    self.logger.info(f"Processing {len(images)} Confluence images")
                    self.stats.total_images += len(images)  # Update images count
                    for idx, image_url in enumerate(images):
                        try:
                            # Add authentication to the image request
                            response = requests.get(
                                image_url,
                                auth=HTTPBasicAuth(
                                    self.confluence_processor.user_name,
                                    self.confluence_processor.api_token
                                ),
                                headers={'Accept': 'image/*'}
                            )

                            if response.status_code == 200:
                                try:
                                    self.logger.info("In Confluence content... trying to open image")
                                    image = Image.open(io.BytesIO(response.content))

                                    # Convert image to RGB if needed
                                    if image.mode in ('RGBA', 'P'):
                                        image = image.convert('RGB')

                                    std_image = image
                                    # std_image = self.file_processor.standardize_image(image)
                                    self.logger.info("In Confluence content... trying to save image")

                                    # Save the image
                                    stored_image_path = self.file_processor.save_image(
                                        response.content,  # Use raw bytes instead of PIL image
                                        f"confluence_{self.config.CONFLUENCE_PAGE_ID}",
                                        0,
                                        idx,
                                        self.config.IMAGE_STORAGE_PATH,
                                        imghdr.what(None, h=response.content) or 'png'  # Detect format from bytes
                                    )

                                    # Get basic metadata
                                    metadata = self.file_processor.get_image_metadata(
                                        response.content,
                                        stored_image_path,
                                        0,
                                        idx
                                    )

                                    # Perform semantic analysis with image bytes
                                    try:
                                        image_analysis = self.semantic_analyzer.analyze_image(response.content)

                                        # Update metadata with semantic analysis
                                        metadata.image_caption = image_analysis.get('image_caption', '')
                                        metadata.image_category = image_analysis.get('image_category', [])
                                        metadata.text_in_image = image_analysis.get('text_in_image', '')

                                    except Exception as e:
                                        self.logger.error(
                                            f"Error in semantic analysis for confluence_{self.config.CONFLUENCE_PAGE_ID}_page0_img{idx}: {str(e)}")
                                        metadata.image_caption = ''
                                        metadata.image_category = []
                                        metadata.text_in_image = ''

                                    # Get embedding and store
                                    embedding = self._get_image_embedding(response.content)

                                    # Add context about the source
                                    related_text = f"Image from Confluence page: {title}\nContext: {text[:500] if text else ''}"

                                    self.vector_store.store_image_vectors(
                                        [embedding],
                                        [metadata],
                                        related_text
                                    )

                                    self.logger.info(
                                        f"Successfully processed Confluence image {idx + 1} from confluence_{self.config.CONFLUENCE_PAGE_ID}")

                                except Exception as e:
                                    self.logger.error(f"Error processing Confluence image {image_url}: {str(e)}")
                                    continue
                            else:
                                self.logger.error(f"Failed to fetch image. Status code: {response.status_code}")
                        except Exception as e:
                            self.logger.error(f"Error processing Confluence image {image_url}: {str(e)}")
                            continue

        except Exception as e:
            self.logger.error(f"Error processing Confluence content: {str(e)}")
            self.stats.errors.append(f"Error processing Confluence content: {str(e)}")


class ImageSemanticAnalyzer:
    def __init__(self, api_key: str, supported_formats: List[str]):
        self.client = Anthropic(api_key=api_key)
        self.model = "claude-3-sonnet-20240229"
        self.supported_formats = [fmt.lower().strip('.') for fmt in supported_formats]
        self.mime_types = {
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'webp': 'image/webp'
        }
        # Set up logging
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)

    def analyze_image(self, image_data: bytes) -> Dict[str, Union[str, List[str]]]:
        """Analyze image using Claude 3.5 Sonnet."""
        try:
            # Set default values
            default_analysis = {
                'image_caption': '',
                'image_category': [],
                'text_in_image': ''
            }
            image_base64 = base64.b64encode(image_data).decode('utf-8')

            response = self.client.messages.create(
                model=self.model,
                max_tokens=500,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": self._get_analysis_prompt()},
                        {"type": "image", "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": image_base64
                        }}
                    ]
                }]
            )

            # Parse response with error handling
            try:
                analysis = json.loads(response.content[0].text)
                return {
                    'image_caption': analysis.get('caption', ''),
                    'image_category': analysis.get('categories', []),
                    'text_in_image': analysis.get('text', '')
                }
            except (json.JSONDecodeError, IndexError, KeyError) as e:
                self.logger.error(f"Error parsing Claude response: {str(e)}")
                return default_analysis

        except Exception as e:
            self.logger.error(f"Error in Claude analysis: {str(e)}")
            return default_analysis

    @staticmethod
    def _get_analysis_prompt() -> str:
        """Get the analysis prompt template."""
        return """Analyze this image and provide:
            1. Short technical caption (1-2 sentences)
            2. Main categories and subcategories
            3. Any text visible in the image
            Return as JSON with these keys: caption, categories, text."""

def setup_logging():
    """Configure logging settings."""
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler('document_processing.log'),
            logging.StreamHandler()
        ]
    )
    return logging.getLogger(__name__)


class ProcessingStats:
    """Track processing statistics."""
    def __init__(self):
        self.total_documents = 0
        self.total_pages = 0
        self.total_images = 0
        self.total_chunks = 0
        self.errors = []
        self.start_time = datetime.now()

    def generate_report(self) -> str:
        """Generate processing report."""
        duration = datetime.now() - self.start_time

        error_details = "\n".join(self.errors) if self.errors else "None"

        return f"""
Processing Report
----------------
Duration: {duration}
Documents Processed: {self.total_documents}
Pages Processed: {self.total_pages}
Images Processed: {self.total_images}
Text Chunks Created: {self.total_chunks}
Average Processing Time per Document: {duration.total_seconds() / max(1, self.total_documents):.2f} seconds
Average Processing Time per Page: {duration.total_seconds() / max(1, self.total_pages):.2f} seconds
Errors Encountered: {len(self.errors)}

Error Details:
{error_details}
"""


def verify_environment(config) -> bool:
    """Verify all necessary environment variables and dependencies."""
    required_attrs = [
        'OPENAI_API_KEY',
        'CONFLUENCE_USERNAME',
        'CONFLUENCE_API_TOKEN',
        'CONFLUENCE_PAGE_ID',
        'DOCUMENT_DIRECTORY',
        'PDF_DIRECTORY',
        'PPT_DIRECTORY',
        'EXCEL_DIRECTORY',
        'LOCAL_QDRANT_PATH',
        'COLLECTION_NAME',
        'CHUNK_SIZE',
        'CHUNK_OVERLAP',
        'BATCH_SIZE',
        'IMAGE_STORAGE_PATH'
    ]
    
    missing_attrs = [attr for attr in required_attrs if not hasattr(config, attr)]
    if missing_attrs:
        raise ValueError(f"Missing required configuration attributes: {missing_attrs}")
    
    return True


def init_directories(config) -> None:
    """Initialize directory structure."""
    directories = [
        config.DOCUMENT_DIRECTORY,
        config.PDF_DIRECTORY,
        config.PPT_DIRECTORY,
        config.EXCEL_DIRECTORY,
        config.LOCAL_QDRANT_PATH,
        config.IMAGE_STORAGE_PATH
    ]

    for directory in directories:
        os.makedirs(directory, exist_ok=True)


def process_documents(config_instance) -> Tuple[bool, ProcessingStats]:
    """
    Process documents and return processing status and statistics.

    Args:
        config_instance: Configuration instance containing all necessary settings

    Returns:
        Tuple[bool, ProcessingStats]: Success status and processing statistics
    """

    logger = setup_logging()
    stats = ProcessingStats()

    try:
        from config import Config
        
        # Verify environment and initialize
        logger.info("Verifying environment configuration...")
        verify_environment(config_instance)
        
        logger.info("Initializing directory structure...")
        init_directories(config_instance)

        # Initialize processor
        logger.info("Initializing document processor...")
        processor = DocumentProcessor(config_instance, stats)

        # Process documents
        logger.info("Starting document processing...")
        num_processed = processor.final_processing()
        stats.total_chunks = num_processed

        # Generate and log report
        report = stats.generate_report()
        logger.info("Processing completed. Final report:\n%s", report)

        return True, stats

    except Exception as e:
        logger.error("Fatal error in main execution: %s", str(e), exc_info=True)
        stats.errors.append(str(e))  # Track errors
        return False, stats


class DocumentProcessingException(Exception):
    """Custom exception for document processing errors."""
    pass


def clean_text(text: str) -> str:
    """Clean and normalize text content."""
    if not text:
        return ""
    
    # Remove excessive whitespace
    text = " ".join(text.split())
    
    # Remove control characters
    text = "".join(char for char in text if char.isprintable())
    
    return text.strip()


def calculate_chunk_size(text: str, target_size: int = 1000) -> int:
    """Calculate optimal chunk size based on text content."""
    words = text.split()
    avg_word_length = sum(len(word) for word in words) / len(words) if words else 5
    return max(100, min(2000, int(target_size / avg_word_length) * 100))


class BatchProcessor:
    """Handle batch processing of documents."""
    def __init__(self, batch_size: int):
        self.batch_size = batch_size
        self.current_batch = []


    def add_item(self, item):
        """Add item to current batch."""
        self.current_batch.append(item)


    def is_ready(self) -> bool:
        """Check if batch is ready for processing."""
        return len(self.current_batch) >= self.batch_size


    def process_batch(self, processor_func):
        """Process current batch and clear it."""
        if self.current_batch:
            processor_func(self.current_batch)
            self.current_batch = []


if __name__ == "__main__":
    # Set up exception handling for the main thread

    from config import Config

    success, stats = process_documents(Config)
    sys.exit(0 if success else 1)